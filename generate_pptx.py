# /// script
# dependencies = [
#   "python-pptx",
#   "google-genai",
# ]
# ///

import os
import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def apply_background(slide, color_rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def add_title(slide, text, color_rgb, size=40, top=1.0):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9.0), Inches(1.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = 'Arial'
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color_rgb
    return title_box

def create_presentation(student_name, lang, slides_data, filepath):
    prs = Presentation()
    
    # Set slide dimensions (16:9 widescreen)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Color Palette: Classic Deep Navy & Coral
    NAVY = RGBColor(28, 40, 51)
    OFF_WHITE = RGBColor(244, 246, 246)
    CORAL = RGBColor(254, 68, 71)
    SLATE = RGBColor(46, 64, 83)
    WHITE = RGBColor(255, 255, 255)

    # Slide 1: Title (Navy background)
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide1 = prs.slides.add_slide(slide_layout)
    apply_background(slide1, NAVY)

    # Title text
    add_title(slide1, "VINSCHOOL MATH EXPLORER" if lang == 'en' else "HỆ THỐNG TOÁN HỌC VINSCHOOL", CORAL, size=38, top=1.2)
    
    sub_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(8.0), Inches(1.5))
    tf = sub_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = f"Class Review & Socratic Feedback Portfolio" if lang == 'en' else "Hồ Sơ Ôn Tập & Gợi Ý Lập Luận Toán Học"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = 'Arial'
    p1.font.size = Pt(20)
    p1.font.color.rgb = WHITE
    p1.font.bold = True
    
    p2 = tf.add_paragraph()
    p2.text = f"Prepared for: {student_name} (Cambridge Stage 6)" if lang == 'en' else f"Biên soạn cho học sinh: {student_name} (Lớp 6A1)"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = 'Arial'
    p2.font.size = Pt(14)
    p2.font.color.rgb = OFF_WHITE
    p2.font.italic = True

    # Slides 2-4 (White background)
    for idx, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(slide_layout)
        apply_background(slide, OFF_WHITE)
        
        # Slide Title
        add_title(slide, slide_data['title'], NAVY, size=28, top=0.5)

        # Content box
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.6))
        tf = content_box.text_frame
        tf.word_wrap = True

        for p_idx, text in enumerate(slide_data['points']):
            p = tf.add_paragraph() if p_idx > 0 else tf.paragraphs[0]
            p.text = f"• {text}"
            p.alignment = PP_ALIGN.LEFT
            p.font.name = 'Arial'
            p.font.size = Pt(14)
            p.font.color.rgb = SLATE
            p.space_after = Pt(12)

            # Highlight specific keywords if necessary
            if ":" in text:
                parts = text.split(":")
                p.text = ""
                run1 = p.add_run()
                run1.text = "• " + parts[0] + ":"
                run1.font.bold = True
                run1.font.color.rgb = NAVY
                
                run2 = p.add_run()
                run2.text = parts[1]
                run2.font.color.rgb = SLATE

    prs.save(filepath)
    print(f"PowerPoint saved successfully to: {filepath}")

def get_personalized_slides(student_id, lang):
    if lang == 'en':
        slides_data = [
            {
                "title": "Concept Breakdown & Gaps",
                "points": [
                    "Target Area: Decimal Place Values (Thập phân) and Vector Translations.",
                    "Common Mistake: Shifting the decimal point by incorrect places when scaling numbers by 10 or 100.",
                    "Socratic Tip: Count the number of zeros in the multiplier. This tells you exactly how many columns to shift the digits."
                ]
            },
            {
                "title": "Interactive Homework Review",
                "points": [
                    "Exercise 1: What value does 9 represent in 14.895? Compare it to 9 in 14.985.",
                    "Exercise 2: Perform translation: Shift a triangle with vertex A(2,3) by vector (-3, 2). Find A'.",
                    "Exercise 3: Write 75% as a simplified fraction. Explain the steps of division."
                ]
            },
            {
                "title": "Pedagogical Guides for Teachers",
                "points": [
                    "Socratic Scaffolding: Ask the student 'What does each column to the right of the decimal represent?'",
                    "Visual Tools: Encourage the student to use the place-value grid or draw coordinate lines for vectors.",
                    "Reflection Goal: Ask the student to write down why place values change when dividing."
                ]
            }
        ]
    else:
        slides_data = [
            {
                "title": "Phân Tích Khái Niệm & Lỗ Hổng",
                "points": [
                    "Nội dung cần củng cố: Giá trị hàng thập phân và Tọa độ Phép biến hình.",
                    "Lỗi thường gặp: Dịch dấu phẩy sai hàng khi nhân chia nhẩm với 10 hoặc 100.",
                    "Gợi ý tư duy: Đếm số chữ số 0 của số nhân. Số chữ số 0 chính là số hàng chữ số cần dịch chuyển."
                ]
            },
            {
                "title": "Bài Tập Rà Soát Tương Tác",
                "points": [
                    "Bài tập 1: Số 9 trong số 14,895 chỉ hàng nào? Nó có giá trị lớn hơn hay nhỏ hơn số 9 trong số 14,985?",
                    "Bài tập 2: Thực hiện phép biến đổi: Tịnh tiến tam giác có đỉnh A(2,3) theo vectơ (-3, 2). Tìm A'.",
                    "Bài tập 3: Viết 75% dưới dạng phân số tối giản và trình bày cách rút gọn."
                ]
            },
            {
                "title": "Gợi Ý Sư Phạm Cho Giáo Viên",
                "points": [
                    "Đặt câu hỏi gợi mở: Hỏi học sinh 'Mỗi hàng bên phải dấu phẩy biểu diễn phân số mấy?'",
                    "Công cụ trực quan: Khuyến khích học sinh dùng lưới tọa độ hoặc bảng dịch hàng số để tính toán.",
                    "Mục tiêu tự ngẫm: Yêu cầu học sinh tự giải thích vì sao giá trị các hàng số giảm đi 10 lần khi dịch sang phải."
                ]
            }
        ]

    # Dynamically prompt Gemini if key exists
    api_key = os.environ.get("GEMINI_API_KEY")
    if api_key and api_key.strip() != "" and "MY_GEMINI_API_KEY" not in api_key:
        try:
            from google import genai
            client = genai.Client(api_key=api_key)
            
            prompt = f"""
            You are a Cambridge Stage 6 Math Curriculum specialist for Vinschool.
            Write a 3-slide classroom review presentation text for student '{student_id}' in language '{lang}'.
            Focus on decimals place values and transformations.
            Return ONLY a valid JSON array of objects matching this schema:
            [
                {{
                    "title": "Slide Title",
                    "points": ["Bullet 1", "Bullet 2", "Bullet 3"]
                }},
                ... (exactly 3 slides)
            ]
            Do not wrap the response in markdown blocks. Output raw JSON string only.
            """
            
            response = client.models.generate_content(
                model='gemini-2.5-flash',
                contents=prompt
            )
            
            data = json.loads(response.text.strip())
            if isinstance(data, list) and len(data) == 3:
                return data
        except Exception as e:
            print(f"Gemini API presentation generation bypassed: {e}. Using official Stage 6 fallbacks.")

    return slides_data

if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: python generate_pptx.py <student_id> <language> <output_path>")
        sys.exit(1)

    student_id = sys.argv[1]
    lang = sys.argv[2]
    output_path = sys.argv[3]
    
    name_map = {
        'student_minh': 'Hoàng Minh',
        'student_nam': 'Lê Nam',
        'student_mai': 'Nguyễn Mai',
        'student_chloe': 'Chloe Smith',
        'student_duc': 'Trần Đức',
        'student_linh': 'Phạm Linh'
      }
    student_name = name_map.get(student_id, "Vinschool Student")

    print(f"Generating review PPTX for {student_name}...")
    slides_data = get_personalized_slides(student_id, lang)
    
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    create_presentation(student_name, lang, slides_data, output_path)
